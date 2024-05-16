import logging
from flask import Flask, request, render_template, redirect, abort, jsonify
from flask_sqlalchemy import SQLAlchemy
from sqlalchemy import Integer, Enum, ForeignKey, Date
from sqlalchemy.exc import SQLAlchemyError
import os
from pptx import Presentation as PPTPresentation
from PyPDF2 import PdfFileReader
from datetime import datetime
from urllib.parse import quote as url_quote
from waitress import serve



# Set up logging
logging.basicConfig(level=logging.DEBUG)

app = Flask(__name__)
basedir = os.path.abspath(os.path.dirname(__file__))
app.config['SQLALCHEMY_DATABASE_URI'] = 'mysql://root:@localhost/db_len'
app.config['UPLOAD_FOLDER'] = os.path.join(basedir, 'static', 'file')
db = SQLAlchemy(app)

class File(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    title = db.Column(db.String(50), nullable=False)
    category = db.Column(db.Enum('Internal', 'External'), nullable=False)
    date = db.Column(db.Date, nullable=False)
    event = db.Column(db.String(25), nullable=False)
    tag = db.Column(db.String(50), nullable=False)
    local_link = db.Column(db.String(255), nullable=False)
    keywords = db.relationship('Keyword', cascade="all,delete", back_populates='file')

class Keyword(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    file_id = db.Column(db.Integer, ForeignKey('file.id'), nullable=False)
    page_number = db.Column(db.Integer, nullable=False)
    keyword = db.Column(db.String(255), nullable=False)
    file = db.relationship("File", back_populates="keywords")

def process_file(path, file_id, category, date, event, tag):
    with app.app_context():
        logging.debug(f"Processing file at path: {path} with file_id: {file_id}")
        title = os.path.basename(path)
        existing_file = File.query.filter_by(id=file_id).first()
        if existing_file:
            logging.debug(f"Updating file with id {file_id}")
            existing_file.category = category
            existing_file.date = date
            existing_file.event = event
            existing_file.tag = tag
            if title:
                existing_file.title = title
                existing_file.local_link = 'file/' + title
            try:
                rows = 0
                if title:
                    rows = Keyword.query.filter_by(file_id=existing_file.id).delete()
                db.session.commit()
                logging.debug(f"Updated file with id {file_id} and removed {rows} previously stored keywords")
            except SQLAlchemyError as e:
                db.session.rollback()
                logging.error(f"Failed updating file: {e}")
        else:
            new_file_entry = File(title=title, category=category, date=date, event=event, tag=tag, local_link='file/' + title)
            db.session.add(new_file_entry)
            db.session.flush()
            logging.debug("New file added to database")
            file_id = new_file_entry.id
        if path.endswith('.pptx'):
            process_presentation(path, file_id, title)
        elif path.endswith('.pdf'):
            process_pdf(path, file_id, title)

def process_presentation(path, file_id, title):
    logging.debug("Processing ppt")
    ppt = PPTPresentation(path)
    existing_entries = Keyword.query.filter_by(file_id=file_id)
    existing_page_numbers = {entry.page_number for entry in existing_entries}
    for i, slide in enumerate(ppt.slides):
        if i + 1 not in existing_page_numbers:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    keyword = shape.text
                    new_entry = Keyword(file_id=file_id, page_number=i + 1, keyword=keyword)
                    db.session.add(new_entry)
    db.session.commit()
    logging.debug(f"Data added to database from presentation: {title}")

def process_pdf(path, file_id, title):
    logging.debug("Processing pdf")
    with open(path, 'rb') as file:
        pdf = PdfFileReader(file)
        existing_entries = Keyword.query.filter_by(file_id=file_id)
        existing_page_numbers = {entry.page_number for entry in existing_entries}
        for i, page in enumerate(pdf.pages):
            if i + 1 not in existing_page_numbers:
                text = page.extract_text()
                for keyword in text.split('\n'):
                    if keyword.strip():
                        new_entry = Keyword(file_id=file_id, page_number=i + 1, keyword=keyword)
                        db.session.add(new_entry)
    db.session.commit()
    logging.debug(f"Data added to database from PDF: {title}")

def process_upload(request):
    file_id = request.form.get('fileId')
    if 'file' not in request.files and file_id is None:
        logging.error("No file and no file_id provided")
        return
    category = request.form['kategori']
    date = datetime.strptime(request.form['tanggal'], '%Y-%m-%d')
    event = request.form['event']
    tag = request.form['tag']
    file = request.files.get("file")
    path = ''
    if file:
        if file.filename == '' and file_id is None:
            logging.error('File is none')
            return
        filename = file.filename
        path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(path)
        logging.debug(f'File is saved: {path}')
    else:
        logging.debug("No file is uploaded")
    process_file(path, int(file_id or -1), category, date, event, tag)

@app.route('/', methods=['GET'])
def index():
    return render_template('index.html')

@app.route('/upload', methods=['GET', 'POST'])
def upload_file():
    if request.method == 'GET':
        return render_template('upload_File.html', file=None, editing=False)
    elif request.method == 'POST':
        process_upload(request)
        return jsonify({'success': True}), 201

@app.route('/edit', methods=['GET', 'POST'])
def edit():
    if request.method == 'GET':
        file_id = request.args.get('id', -1)
        if file_id == -1:
            abort(404)
        file = File.query.filter_by(id=file_id).first_or_404()
        return render_template('upload_File.html', file=file, editing=True)
    elif request.method == 'POST':
        process_upload(request)
        return jsonify({'success': True}), 200

@app.route('/delete', methods=['GET'])
def delete():
    file_id = request.args.get('id', -1)
    if file_id == -1:
        logging.error('ID file tidak ada, gagal menghapus')
        return redirect(request.referrer)
    keyword_count = Keyword.query.filter_by(file_id=file_id).delete()
    count = File.query.filter_by(id=file_id).delete()
    db.session.commit()
    logging.debug(f'Menghapus {count} file dengan id {file_id} dan {keyword_count} keyword')
    return redirect(request.referrer)

@app.route('/search', methods=['GET'])
def search():
    page = int(request.args.get('page', 1))
    mode = request.args.get('mode')
    search_keyword = request.args.get('keyword', "")
    query = None
    if mode == 'files':
        contains_keyword = File.query.filter(File.tag.like(f"%{search_keyword}%"))
        query = contains_keyword.paginate(page=page, per_page=10)
    else:
        contains_keyword = Keyword.query.join(Keyword.file).filter(Keyword.keyword.like(f"%{search_keyword}%")).group_by(
            File.title, Keyword.page_number)
        query = contains_keyword.paginate(page=page, per_page=10)
    return render_template('results.html', results=query, keyword=search_keyword, mode=mode)

@app.route('/files', methods=['GET'])
def files():
    sort = request.args.get('sort', "desc")
    page = int(request.args.get('page', 1))
    paginated_files = File.query.order_by(File.date.desc() if sort == "desc" else File.date.asc()).paginate(page=page, per_page=10)
    return render_template('files.html', paginated_files=paginated_files, sort=sort)

if __name__ == '__main__':
    from waitress import serve
    with app.app_context():
        db.create_all()
    logging.debug("Starting the server")
    serve(app, host='127.0.0.1', port=5000)
